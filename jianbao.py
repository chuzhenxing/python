# coding:utf-8
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs
########################################################################################################################
font_size = 8  # 字体大小
fig_size = (10, 4)  # 图表大小
bar_width = 0.35   # 设置柱形图宽度
########################################################################################################################
plt.rcParams['font.size'] = font_size   # 更改默认更新字体大小
plt.rcParams['figure.figsize'] = fig_size   # 修改默认更新图表大小
plt.rcParams['font.sans-serif'] = ['Microsoft Yahei']  # 指定默认字体
plt.rcParams['axes.unicode_minus'] = False  # 解决保存图像是负号'-'显示为方块的问题
# ################# 柱状图函数汇总 #####################################################################################
# 添加数据标签 就是矩形上面的数字
def add_labels(recs, gaodu):
    for rect in recs:
        height = int(rect.get_height())
        plt.text(rect.get_x() + rect.get_width() / 2, height/gaodu, height, ha='center', va='bottom',color='k',fontsize=8)
        # horizontalalignment='center' plt.text(x坐标，y坐标，text,位置)
        # 柱形图边缘用白色填充，为了更加清晰可分辨
        rect.set_edgecolor('white')
def add_labels5(recs, gaodu):
    for rect in recs:
        height = int(rect.get_height())
        plt.text(rect.get_x() + rect.get_width() / 2, rect.get_height()/gaodu, rect.get_height(), ha='center', va='bottom',color='k',fontsize=8)
        # horizontalalignment='center' plt.text(x坐标，y坐标，text,位置)
        # 柱形图边缘用白色填充，为了更加清晰可分辨
        rect.set_edgecolor('white')
def add_labels2(recs1, recs2):    # 2层叠加图
    for rec11, rec22 in zip(recs1,recs2):
        height11 = int(rec11.get_height())
        height22 = int(rec22.get_height())
        plt.text(rec11.get_x() + rec11.get_width() / 2, height11+ height22*2/5, height22, ha='center', va='bottom',color='k',fontsize=8)
        rec22.set_edgecolor('white')
def add_labels3(recs, k):         # 6层叠加图
    for rect,temp in zip(recs, k):
        height = int(rect.get_height())
        plt.text(rect.get_x() + rect.get_width() / 2, temp+height*2/5, height, ha='center', va='bottom',color='k',fontsize=8)
        # horizontalalignment='center' plt.text(x坐标，y坐标，text,位置)
        # 柱形图边缘用白色填充，为了更加清晰可分辨
        rect.set_edgecolor('white')
# #################线状图函数汇总 ######################################################################################
def add_labels4(index_1, scores_1,scores_2):
    for a,b,c in zip(index_1 + bar_width, scores_1, scores_2):
        plt.text(a,b ,c,ha = 'center',va = 'bottom',fontsize = 8, color='k')

# ##################导入excel###########################################################################################
# 1.导入定时任务excel
# df = pd.read_excel('出借周一汇报_12-04.xlsx'.decode('utf-8')  , sheetname = range(0,21))
df = pd.read_excel('出借周一汇报_12-04.xlsx'.decode('utf-8'), sheetname = [u'出借简报_目标达成概况（月累计）1', u'出借简报_出借周趋势2', u'出借简报_新增出借人数周趋势3', u'出借简报_出借金额分产品4',
                                                                       u'出借简报_出借转化率6',u'出借简报_二次出借转化（新手标）7',u'出借简报_AUM时点月走势8',u'出借简报_用户分层9',u'出借简报_一起赚10_1',
                                                                       u'出借简报_一起赚10_2',u'出借简报_一起赚10_3',u'出借简报_一起赚10_4',u'出借简报_一起赚10_5',u'出借简报_积分11_1',u'出借简报_积分11_2',
                                                                       u'出借简报_商城12_1',u'出借简报_商城分产品12_2',u'出借简报_基金13_1',u'出借简报_基金13_2',u'出借简报_保险14'])
# 2.导入模板
dfbiao = pd.read_excel('模板_出借业务分渠道统计.xlsx'.decode('utf-8'), sheetname = [u'周报'],  skiprows = 7, skip_footer = 15)

# #############################################################核心指标-目标达成概况#######################ppt第三页图片

plt.figure(figsize=(10, 4))
# plt.tight_layout()
plt.subplots_adjust(left=0.07, bottom=0.1, right=0.95, top=0.9, hspace=0.8, wspace=0.3)   #调整图像边缘及图像间的空白间隔

ax1 = plt.subplot(121)  # 在图表1中创建子图1
ax2 = plt.subplot(122)  # 在图表1中创建子图2

plt.sca(ax1)    # 选择图表1的子图1
# df2：出借简报_出借周趋势2 数据框
df2 = df[u'出借简报_出借周趋势2'].round( 0)
df2[u'出借金额'] = df2[u'出借金额'].astype('int64')
df2[u'人均出借金额'] = df2[u'人均出借金额'].astype('int64')
# print df2
# print df2.dtypes
timespan = tuple(df2[u'日期'])     # 时段段元组
# print timespan
names = (u'出借金额', u'人均出借金额',u'出借人数')  # 名称元组
# print names
scores = (tuple(df2[u'出借金额']), tuple(df2[u'人均出借金额']), tuple(df2[u'出借人数']))  # 金额元组
# print scores

df_total11 = df2['出借金额'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围

index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax1.bar(index, scores[0], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
rec2 = ax1.bar(index + bar_width, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'金额（万元）',fontsize=10)     #y轴标签
plt.title(u'出借周趋势',fontsize=10)     # 子图1标题

plt.legend(loc = 'upper center',bbox_to_anchor = (0.38,1.0),fancybox = True,ncol = 3,fontsize = 12.5,frameon = False,prop={'size':8})

add_labels(rec1, 1)    # 画柱状图标签
add_labels(rec2, 1)

plt.twinx()               #副坐标
df_total12 = df2['出借人数'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*2, ymin=0)        # Y轴范围

plt.ylabel(u'出借人数（人）',fontsize=10)
plt.plot(index + bar_width, scores[2] , label=names[2],linewidth=2, color = 'orange')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.8,1.0),fancybox = True,ncol = 1,fontsize = 8,frameon = False,prop={'size':8})

add_labels4(index, scores[2],scores[2])   # 画折线图标签


plt.xticks(index + bar_width, timespan)        # X轴标题

plt.savefig(u'出借简报_出借周趋势2.png')   # 图表输出到本地
# pylab.show('scores_par.png')    # 并打印显示图片
# plt.show()

#################
plt.sca(ax2)    # 选择图表1的子图2
# df3：出借简报_新增出借人数周趋势3 数据框
df3 = df[u'出借简报_新增出借人数周趋势3'].round(decimals = 0)
# print df3.dtypes
timespan = tuple(df3[u'日期'])     # 时段段元组
# print timespan
names = (u'新注册新出借用户', u'老注册新出借用户')  # 名称元组
# print names
scores = (tuple(df3[u'新注册新出借用户']), tuple(df3[u'老注册新出借用户']))  # 金额元组
# print scores

df_total11 = df3['新注册新出借用户'.decode('utf-8')].max()+df3['老注册新出借用户'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围

index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax2.bar(index, scores[0], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
rec2 = ax2.bar(index, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1],bottom=scores[0])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
# plt.ylabel(u'新老出借用户数',fontsize=10)     # y轴标签
plt.title(u'新增出借人数周趋势',fontsize=10)     # 子图1标题

plt.legend(loc = 'upper center',bbox_to_anchor = (0.5,1.0),fancybox = True,ncol = 2,fontsize = 12.5,frameon = False,prop={'size':8})
plt.xticks(index + bar_width/2, timespan)        # X轴标题

add_labels(rec1, 2)     # 画柱状图标签
add_labels2(rec1, rec2)   # 画柱状图标签
plt.savefig(u'出借简报_新增出借人数周趋势3.png')   # 图表输出到本地
# pylab.show('scores_par.png')    # 并打印显示图片
# plt.show( )

###############################################出借简报_出借金额分产品4#########ppt第5页####################
plt.figure(figsize=(10, 5))
# plt.tight_layout()
plt.subplots_adjust(left=0.1, bottom=0.1, right=0.9, top=0.9,hspace=0.8, wspace=0.4)   #调整图像边缘及图像间的空白间隔
ax3 = plt.subplot(111)  # 在图表1中创建子图1

# df4：出借简报_出借金额分产品4 数据框
df4 = df[u'出借简报_出借金额分产品4'].round(decimals = 0)
df4[u'一月期_出借金额'] = df4[u'一月期_出借金额'].astype('int64')
df4[u'三月期_出借金额'] = df4[u'三月期_出借金额'].astype('int64')
df4[u'六月期_出借金额'] = df4[u'六月期_出借金额'].astype('int64')
df4[u'九月期_出借金额'] = df4[u'九月期_出借金额'].astype('int64')
df4[u'十二月期_出借金额'] = df4[u'十二月期_出借金额'].astype('int64')
df4[u'二十四月期_出借金额'] = df4[u'二十四月期_出借金额'].astype('int64')
# print df4.dtypes
timespan = tuple(df4[u'日期'])     # 时段段元组
# print timespan
names = (u'一月期', u'三月期', u'六月期', u'九月期' , u'十二月期' ,u'二十四月期' , u'平均产品周期')  # 名称元组
# print names
scores = (tuple(df4[u'一月期_出借金额']), tuple(df4[u'三月期_出借金额']), tuple(df4[u'六月期_出借金额']),tuple(df4[u'九月期_出借金额']), tuple(df4[u'十二月期_出借金额']), tuple(df4[u'二十四月期_出借金额']), tuple(df4[u'平均产品周期']))  # 金额元组
# print scores
gg=[]
for a,b,c,d,e,f in zip(df4[u'一月期_出借金额'],df4[u'三月期_出借金额'],df4[u'六月期_出借金额'],df4[u'九月期_出借金额'],df4[u'十二月期_出借金额'],df4[u'二十四月期_出借金额']):
    gg.append(a+b+c+d+e+f)
# print gg
# print max(gg)
df_total44 = max(gg)
plt.ylim(ymax=df_total44*1.5, ymin=0)        # Y轴范围
index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# print scores[0]
# print scores[1]
# print scores[2]
temp0=tuple(pd.Series(list(scores[0])))
temp1=tuple(pd.Series(list(scores[0])) + pd.Series(list(scores[1])))
temp2=tuple(pd.Series(list(scores[0])) + pd.Series(list(scores[1]))+ pd.Series(list(scores[2])))
temp3=tuple(pd.Series(list(scores[0])) + pd.Series(list(scores[1]))+ pd.Series(list(scores[2]))+ pd.Series(list(scores[3])))
temp4=tuple(pd.Series(list(scores[0])) + pd.Series(list(scores[1]))+ pd.Series(list(scores[2]))+ pd.Series(list(scores[3]))+ pd.Series(list(scores[4])))
# print temp2
# temp2=tuple(pd.DataFrame(list(scores[0])) + pd.DataFrame(list(scores[1])))
# print temp1,temp2
# # matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax3.bar(index, scores[0], bar_width, linewidth=0, color='chocolate', label=names[0])     # 第一个柱状图
rec2 = ax3.bar(index, scores[1], bar_width, linewidth=0,  color='sandybrown', label=names[1],bottom=temp0)   # 第二个柱状图
rec3 = ax3.bar(index, scores[2], bar_width, linewidth=0, color='bisque', label=names[2],bottom=temp1)     # 第三个柱状图
rec4 = ax3.bar(index, scores[3], bar_width, linewidth=0,  color='gold', label=names[3],bottom=temp2)   # 第四个柱状图
rec5 = ax3.bar(index, scores[4], bar_width, linewidth=0, color='deepskyblue', label=names[4],bottom=temp3)     # 第五个柱状图
rec6 = ax3.bar(index, scores[5], bar_width, linewidth=0,  color='steelblue', label=names[5],bottom=temp4)   # 第六个柱状图

plt.xlabel(u'日期（近六周）',fontsize=10)    # x轴标签
plt.ylabel(u'金额（万元）',fontsize=10)     # y轴标签
plt.title(u'出借金额（万元）',fontsize=10)     # 标题
plt.legend(loc = 'upper center',fancybox = True,ncol = 7,fontsize = 8,frameon = False,prop={'size':8})
plt.xticks(index + bar_width/2, timespan)        # X轴标题

add_labels(rec1, 2.5)               # 画柱状图标签
add_labels3(rec2, temp0)             # 画柱状图标签
add_labels3(rec3, temp1)
add_labels3(rec4, temp2)
add_labels3(rec5, temp3)
add_labels3(rec6, temp4)


plt.twinx()               #副坐标
df_total12 = df4['平均产品周期'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*1.5, ymin=0)        # Y轴范围

plt.ylabel(u'平均产品周期',fontsize=10)
plt.plot(index + bar_width/2, scores[6] , label=names[6],linewidth=2, color = 'orange')    #画折线图
# plt.legend(loc = 'upper center',bbox_to_anchor = (0.8,1.0),fancybox = True,ncol = 1,fontsize = 12.5,frameon = False,prop={'size':8})
add_labels4(index, scores[6], scores[6])   # 画折线图标签
plt.savefig(u'出借简报_出借金额分产品4.png')   # 图表输出到本地
# pylab.show('scores_par.png')    # 并打印显示图片
# plt.show( )
########################################出借简报_出借转化率6############ppt第7页#######################################
plt.figure(figsize=(10, 5))
# plt.tight_layout()
plt.subplots_adjust(left=0.1, bottom=0.1, right=0.9, top=0.9,hspace=0.8, wspace=0.4)   #调整图像边缘及图像间的空白间隔
ax5 = plt.subplot(111)  # 在图表1中创建子图1

# df4：出借简报_出借金额分产品4 数据框
df5 = df[u'出借简报_出借转化率6']
# print df5[u'当周出借转化率']
df5[u'当周出借转化率2'] =df5[u'当周出借转化率'].apply(lambda x:format(x,'.0%'))
# print df5.dtypes
# print df5.dtypes
timespan = tuple(df5[u'日期'])     # 时段段元组
# print timespan
names = (u'当周注册人数', u'当周注册且出借人数', u'当周出借转化率')  # 名称元组
# print names
scores = (tuple(df5[u'当周注册人数']), tuple(df5[u'当周注册且出借人数']), tuple(df5[u'当周出借转化率']),tuple(df5[u'当周出借转化率2']))  # 金额元组
# print scores

df_total11 = df5['当周注册人数'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围
index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax5.bar(index, scores[0], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
rec2 = ax5.bar(index + bar_width, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'注册&新出借人数',fontsize=10)     #y轴标签
plt.title(u'出借转化率',fontsize=10)     # 子图1标题

plt.legend(loc = 'upper center',bbox_to_anchor = (0.4,1.0),fancybox = True,ncol = 3,fontsize = 12.5,frameon = False,prop={'size':8})
add_labels(rec1, 1)     # 画柱状图标签
add_labels(rec2, 1)     # 画柱状图标签

plt.twinx()               #副坐标
df_total12 = df5['当周出借转化率'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*2, ymin=0)        # Y轴范围

plt.ylabel(u'当周出借转化率',fontsize=10)
plt.plot(index + bar_width, scores[2] , label=names[2],linewidth=2, color = 'orange')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.65,1.0),fancybox = True,ncol = 1,fontsize = 8,frameon = False,prop={'size':8})

# for a,b,c in zip(index + bar_width, scores[2],scores[3]):   #for循环，将index3 + bar_width / 2、zengzhanglv[0]、zengzhanglvxiabiao[0]分别赋值给参数a、b、c
#     plt.text(a,b ,c,ha = 'center',va = 'bottom',fontsize = 8, color='orange')
add_labels4(index, scores[2],scores[3])   # 画折线图标签

plt.xticks(index + bar_width, timespan)        # X轴标题

plt.savefig(u'出借简报_出借转化率6.png',transparent=True)  # 图表输出到本地
# plt.savefig(u'出借简报_出借转化率6.png', format='png', bbox_inches='tight', transparent=True, dpi=600)  # 图表输出到本地
# plt.show()

#########################################出借简报_二次出借转化（新手标）7############ppt第8页###########################
plt.figure(figsize=(10, 5))
# plt.tight_layout()
plt.subplots_adjust(left=0.1, bottom=0.1, right=0.9, top=0.9,hspace=0.8, wspace=0.4)   #调整图像边缘及图像间的空白间隔
ax5 = plt.subplot(111)  # 在图表1中创建子图1

# df5：出借简报_二次出借转化（新手标）7 数据框
df5 = df[u'出借简报_二次出借转化（新手标）7']
print df5
df5[u'近30天二次出借转化率2'] =df5[u'近30天二次出借转化率'].apply(lambda x:format(x,'.1%'))
print df5.dtypes
# print df5.dtypes
timespan = tuple(df5[u'日期'])     # 时段段元组
# print timespan
names = (u'近30天新出借人数', u'近30天二次出借人数', u'近30天二次出借转化率')  # 名称元组
# print names
scores = (tuple(df5[u'近30天新出借人数']), tuple(df5[u'近30天二次出借人数']), tuple(df5[u'近30天二次出借转化率']),tuple(df5[u'近30天二次出借转化率2']))  # 金额元组
# print scores
df_total11 = df5['近30天新出借人数'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围
index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax5.bar(index, scores[0], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
rec2 = ax5.bar(index + bar_width, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'首借&二借人数',fontsize=10)     #y轴标签
plt.title(u'新手标二次出借转化率（近30天）',fontsize=10)     # 子图1标题

plt.legend(loc = 'upper center',bbox_to_anchor = (0.4,1.0),fancybox = True,ncol = 3,fontsize = 12.5,frameon = False,prop={'size':8})
add_labels(rec1, 1)     # 画柱状图标签
add_labels(rec2, 1)     # 画柱状图标签

plt.twinx()               #副坐标
df_total12 = df5['近30天二次出借转化率'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*2, ymin=0)        # Y轴范围

plt.ylabel(u'近30天二次出借转化率',fontsize=10)
plt.plot(index + bar_width, scores[2] , label=names[2],linewidth=2, color = 'orange')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.65,1.0),fancybox = True,ncol = 1,fontsize = 8,frameon = False,prop={'size':8})

# for a,b,c in zip(index + bar_width, scores[2],scores[3]):   #for循环，将index3 + bar_width / 2、zengzhanglv[0]、zengzhanglvxiabiao[0]分别赋值给参数a、b、c
#     plt.text(a,b ,c,ha = 'center',va = 'bottom',fontsize = 8, color='orange')
add_labels4(index, scores[2],scores[3])   # 画折线图标签

plt.xticks(index + bar_width, timespan)        # X轴标题

plt.savefig(u'出借简报_二次出借转化（新手标）7.png',transparent=True)  # 图表输出到本地
# plt.savefig(u'出借简报_出借转化率6.png', format='png', bbox_inches='tight', transparent=True, dpi=600)  # 图表输出到本地
# plt.show()












########################################  出借简报_AUM时点月走势8 #################ppt第11页############################
plt.figure(figsize=(10, 5.4))
# plt.tight_layout()
plt.subplots_adjust(left=0.1, bottom=0.1, right=0.9, top=0.9, hspace=0.3, wspace=0.3)   #调整图像边缘及图像间的空白间隔

ax1 = plt.subplot(211)  # 在图表1中创建子图1
ax2 = plt.subplot(212)  # 在图表1中创建子图2

plt.sca(ax1)    # 选择图表1的子图1
# df2：出借简报_出借周趋势2 数据框
df2 = df[u'出借简报_AUM时点月走势8'].round(2)
# df2[u'出借金额'] = df2[u'出借金额'].astype('float64')
# df2[u'人均出借金额'] = df2[u'人均出借金额'].astype('float64')
print df2
print df2.dtypes
timespan = tuple(df2['date'])     # 时段段元组
# print timespan
names = (u'在投用户(万人)', u'在投金额(亿元)', u'人均在投金额(万元)')  # 名称元组
# print names
scores = (tuple(df2[u'在投用户万人']), tuple(df2[u'在投金额亿元']), tuple(df2[u'人均在投金额万元']))  # 金额元组
# print scores
df_total11 = df2['在投金额亿元'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围
index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# # matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax1.bar(index, scores[1], bar_width, linewidth=0, color='deepskyblue', label=names[1])     #第一个柱状图
# rec2 = ax1.bar(index + bar_width, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1])   #第二个柱状图
# plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'在投金额（亿元）',fontsize=8)     #y轴标签
plt.title(u'AUM-时点逐月走势',fontsize=10)     # 子图1标题
plt.legend(loc = 'upper center',bbox_to_anchor = (0.38,1.0),fancybox = True,ncol = 3,fontsize = 12.5,frameon = False,prop={'size':8})
#
add_labels(rec1, 1)    # 画柱状图标签
plt.twinx()               #副坐标
df_total12 = df2['在投用户万人'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*3, ymin=0)        # Y轴范围
plt.ylabel(u'在投用户(万人)',fontsize=10)
plt.plot(index + bar_width, scores[0] , label=names[0],linewidth=2, color = 'orange')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.6,1.0),fancybox = True,ncol = 1,fontsize = 8,frameon = False,prop={'size':8})
#
add_labels4(index, scores[0],scores[0])   # 画折线图标签
#
#
plt.xticks(index + bar_width, timespan)        # X轴标题
#
#################
plt.sca(ax2)    # 选择图表1的子图2
df_total11 = df2['人均在投金额万元'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围

index = np.arange(len(scores[2]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax2.bar(index, scores[2], bar_width, linewidth=0, color='deepskyblue', label=names[2])     #第一个柱状图
#  rec2 = ax2.bar(index, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1],bottom=scores[0])   #第二个柱状图
# plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'人均在投金额（万元）',fontsize=10)     # y轴标签
plt.title(u'AUM-时点逐月走势',fontsize=10)     # 子图1标题

plt.legend(loc = 'upper center',bbox_to_anchor = (0.5,1.0),fancybox = True,ncol = 1,fontsize = 12.5,frameon = False,prop={'size':8})
plt.xticks(index + bar_width/2, timespan)        # X轴标题

add_labels5(rec1, 1)     # 画柱状图标签
plt.savefig(u'出借简报_AUM时点月走势8.png',transparent=True)   # 图表输出到本地
# pylab.show('scores_par.png')    # 并打印显示图片
# plt.show()

########################################  出借简报_一起赚10_1和 出借简报_一起赚10_2 #################ppt第16页##########
plt.figure(figsize=(10, 6))
# plt.tight_layout()
plt.subplots_adjust(left=0.07, bottom=0.1, right=0.95, top=0.9, hspace=0.2, wspace=0.2)   #调整图像边缘及图像间的空白间隔
ax1 = plt.subplot(121)  # 在图表1中创建子图1
ax2 = plt.subplot(122)  # 在图表1中创建子图2
plt.sca(ax1)    # 选择图表1的子图1

df1 = df[u'出借简报_一起赚10_1'].round(4)
df2 = df[u'出借简报_一起赚10_2'].round(4)
df3 = df[u'出借简报_一起赚10_3'].round(0)
df4 = df[u'出借简报_一起赚10_4'].round(0)
df5 = df[u'出借简报_一起赚10_5'].round(0)
df1[u'一起赚注册占比互联网2'] =df1[u'一起赚注册占比互联网'].apply(lambda x:format(x,'.1%'))
df2[u'一起赚出借占比互联网2'] =df2[u'一起赚出借占比互联网'].apply(lambda x:format(x,'.1%'))
df2[u'一起赚出借占比互联网2'] =df2[u'一起赚出借占比互联网'].apply(lambda x:format(x,'.1%'))
# print df2
# print df2.dtypes
timespan = tuple(df1[u'日期'])     # 时段段元组
# print timespan
names = (u'一起赚注册人数', u'一起赚注册占比互联网',u'一起赚出借人数',u'一起赚出借占比互联网',u'一起赚出借金额', u'一起赚金额占比互联网')  # 名称元组
# print names

scores = (tuple(df1[u'一起赚注册人数']), tuple(df1[u'一起赚注册占比互联网']), tuple(df2[u'一起赚出借人数']),tuple(df2[u'一起赚出借占比互联网']), tuple(df2[u'一起赚出借金额']), tuple(df2[u'一起赚金额占比互联网']),
          tuple(df1[u'一起赚注册占比互联网2']), tuple(df2[u'一起赚出借占比互联网2']), tuple(df2[u'一起赚出借占比互联网2']))  # 金额元组
# print scores
df_total11 = df2['一起赚出借人数'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围
index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax1.bar(index, scores[0], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
rec2 = ax1.bar(index + bar_width, scores[2], bar_width, linewidth=0,  color='paleturquoise', label=names[2])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
plt.ylabel(u'一起赚注册&出借人数',fontsize=10)     #y轴标签
plt.title(u'一起赚人数',fontsize=10)     # 子图1标题
plt.legend(loc = 'upper center',bbox_to_anchor = (0.48,1.0),fancybox = True,ncol = 2,fontsize = 8,frameon = False,prop={'size':8})
add_labels(rec1, 1)    # 画柱状图标签
add_labels(rec2, 1)

plt.twinx()               #副坐标
df_total12 = df2['一起赚出借占比互联网'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*2, ymin=0)        # Y轴范围

plt.ylabel(u'一起赚人数占比（人）',fontsize=10)
plt.plot(index + bar_width, scores[1] , label=names[1],linewidth=2, color = 'orange')    #画折线图
plt.plot(index + bar_width, scores[3] , label=names[3],linewidth=2, color = 'bisque')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.48,0.95),fancybox = True,ncol = 2,fontsize = 8,frameon = False,prop={'size':8})
add_labels4(index, scores[1],scores[6])   # 画折线图标签
add_labels4(index, scores[3],scores[7])   # 画折线图标签
plt.xticks(index + bar_width, timespan)        # X轴标题

#################
plt.sca(ax2)    # 选择图表1的子图2
df_total11 = df2['一起赚出借金额'.decode('utf-8')].max()
plt.ylim(ymax=df_total11*1.5, ymin=0)        # Y轴范围
#index = np.arange(len(scores[0]))   # 绘制出借金额 index表示柱形图左边x的坐标
# # matplotlib.pyplot.bar(left, height, width=0.8, bottom=None, hold=None, data=None, **kwargs)
rec1 = ax2.bar(index, scores[4], bar_width, linewidth=0, color='deepskyblue', label=names[0])     #第一个柱状图
# rec2 = ax2.bar(index, scores[1], bar_width, linewidth=0,  color='paleturquoise', label=names[1],bottom=scores[0])   #第二个柱状图
plt.xlabel(u'日期（近六周）',fontsize=10)    #x轴标签
# plt.ylabel(u'出借金额（万）',fontsize=10)     # y轴标签
plt.title(u'一起赚出借金额',fontsize=10)     # 子图1标题
plt.legend(loc = 'upper center',bbox_to_anchor = (0.5,1.0),fancybox = True,ncol = 2,fontsize = 12.5,frameon = False,prop={'size':8})
plt.xticks(index + bar_width/2, timespan)        # X轴标题
add_labels(rec1, 1)     # 画柱状图标签


plt.twinx()               #副坐标
df_total12 = df2['一起赚金额占比互联网'.decode('utf-8')].max()
plt.ylim(ymax=df_total12*2, ymin=0)        # Y轴范围

# plt.ylabel(u'一起赚人数占比（人）',fontsize=10)
plt.plot(index + bar_width, scores[5] , label=names[5],linewidth=2, color = 'orange')    #画折线图
plt.legend(loc = 'upper center',bbox_to_anchor = (0.48,0.95),fancybox = True,ncol = 2,fontsize = 8,frameon = False,prop={'size':8})
add_labels4(index, scores[5],scores[8])   # 画折线图标签
plt.xticks(index + bar_width, timespan)        # X轴标题

plt.savefig(u'出借简报_一起赚10_2.png',transparent=True)   # 图表输出到本地
# pylab.show('scores_par.png')    # 并打印显示图片
plt.show( )













######################################################################################################################### 输出到ppt
prs = Presentation(u'宜人理财业务分析周报简版-2017.1.2.pptx')
slide2 = prs.slides[2]        #  第3页幻灯片
left =Inches(0.0)
top = Inches(3.0)
pic = slide2.shapes.add_picture(u'出借简报_新增出借人数周趋势3.png', left, top)
###################################
slide4 = prs.slides[4]        #  第5页幻灯片
left =Inches(0.0)
top = Inches(2.0)
pic = slide4.shapes.add_picture(u'出借简报_出借金额分产品4.png', left, top)
##################################

slide6 = prs.slides[6]        #  第7页幻灯片
left =Inches(0.0)
top = Inches(2.0)
pic = slide6.shapes.add_picture(u'出借简报_出借转化率6.png', left, top)

##################################

slide6 = prs.slides[7]        #  第8页幻灯片
left =Inches(0.0)
top = Inches(2.0)
pic = slide6.shapes.add_picture(u'出借简报_二次出借转化（新手标）7.png', left, top)

##################################

slide10 = prs.slides[10]        #  第11页幻灯片
left =Inches(0.0)
top = Inches(1.4)
pic = slide10.shapes.add_picture(u'出借简报_AUM时点月走势8.png', left, top)

##################################
slide15 = prs.slides[15]        #  第16页幻灯片,一起赚
left =Inches(0.0)
top = Inches(1.4)
pic = slide15.shapes.add_picture(u'出借简报_一起赚10_2.png', left, top)
##################################
slide16 = prs.slides[16]        #  第17页幻灯片,一起赚
left =Inches(0.0)
top = Inches(1.4)
# pic = slide16.shapes.add_picture(u'出借简报_一起赚10_5.png', left, top)



#################################
dfbiao1=dfbiao [u'周报']
# print dfbiao1
dfbiao1_1=dfbiao1.iloc[0:16,1:5]     #截取'模板_出借业务分渠道统计.xlsx'—周报表格切片
dfbiao1_1[u'累计出借金额'] = dfbiao1_1[u'累计出借金额'].astype('float64')
dfbiao1_1[u'总目标'] = dfbiao1_1[u'总目标'].astype('int64')
dfbiao1_1[u'完成率'] = dfbiao1_1[u'完成率'].astype('float64')
# print dfbiao1_1.dtypes
dfbiao1_1[u'完成率'] =dfbiao1_1[u'完成率'].apply(lambda x:format(x,'.2%'))
# print dfbiao1_1
# for r in table.rows:
#         for c in r.cells:
#             s += c.text_frame.text + " | "
#         #to write
#         #c.text_frame.text = "example"
#             print s
slide5 = prs.slides[5]        #  第6页幻灯片,修改表格  核心指标-分渠道构成
table = slide5.shapes[3].table # maybe 0..n
for i in range(2, 5):
    for j in range(2, 18):
        table.cell(j, i).text = str(dfbiao1_1.iloc[j-2, i-1])
################################# 用户分层-月度对比-人数与金额占比
# dffenceng：出借简报_用户分层9  数据框
dffenceng = df[u'出借简报_用户分层9'].round(4)
dffenceng[u'前2月金额'] =dffenceng[u'前2月金额'].astype('int64')
dffenceng[u'前1月金额'] =dffenceng[u'前1月金额'].astype('int64')
dffenceng[u'前0月金额'] =dffenceng[u'前0月金额'].astype('int64')
dffenceng[u'人数占比2_'] =dffenceng[u'人数占比2'].apply(lambda x:format(x,'.2%'))
dffenceng[u'人数占比1_'] =dffenceng[u'人数占比1'].apply(lambda x:format(x,'.2%'))
dffenceng[u'人数占比0_'] =dffenceng[u'人数占比0'].apply(lambda x:format(x,'.2%'))
dffenceng[u'金额占比2_'] =dffenceng[u'金额占比2'].apply(lambda x:format(x,'.2%'))
dffenceng[u'金额占比1_'] =dffenceng[u'金额占比1'].apply(lambda x:format(x,'.2%'))
dffenceng[u'金额占比0_'] =dffenceng[u'金额占比0'].apply(lambda x:format(x,'.2%'))


# print dffenceng
# print dffenceng.dtypes

slide11 = prs.slides[11]        #  第12页幻灯片,修改表格  用户分层-月度对比-人数与金额占比
table = slide11.shapes[2].table # maybe 0..n
for i in range(1, 13):
    for j in range(3, 16):
            table.cell(j, i).text = str(dffenceng .iloc[j-3, i])
            table.cell(j, i).text_frame.paragraphs[0].font.size  = Pt(10)   #   字体大小统一
            if i in [2,3,5,6,8,9,11,12] and  dffenceng .iloc[j-3, i-1] < dffenceng .iloc[j-3, i] :  table.cell(j, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 139, 0)   #  颜色设置绿色
            if i in [2,3,5,6,8,9,11,12] and  dffenceng .iloc[j-3, i-1] > dffenceng .iloc[j-3, i] :  table.cell(j, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)   #  颜色设置红色
            if i in [4,5,6] :  table.cell(j, i).text = str(dffenceng .iloc[j-3, i+9])   # 转变为百分比格式
            if i in [10,11,12] :  table.cell(j, i).text = str(dffenceng .iloc[j-3, i+6])  # 转变为百分比格式
#################################################################

prs.save('new-file-name.pptx')
